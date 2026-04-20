from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
import textwrap

from mdconvert_app.models import ConversionResult


@dataclass
class HeadingBlock:
    level: int
    text: str


@dataclass
class ParagraphBlock:
    text: str


@dataclass
class ListBlock:
    items: list[str]


@dataclass
class TableBlock:
    rows: list[list[str]]


MarkdownBlock = HeadingBlock | ParagraphBlock | ListBlock | TableBlock


def convert_markdown(source: Path, destination: Path) -> ConversionResult:
    markdown = source.read_text(encoding="utf-8")
    blocks = parse_markdown(markdown)
    destination.parent.mkdir(parents=True, exist_ok=True)
    warnings: list[str] = []

    suffix = destination.suffix.lower()
    if suffix == ".docx":
        _write_docx(destination, blocks)
    elif suffix == ".pptx":
        _write_pptx(destination, blocks)
    elif suffix == ".xlsx":
        _write_xlsx(destination, blocks)
    elif suffix == ".pdf":
        _write_pdf(destination, blocks)
    else:
        raise ValueError(f"Unsupported Markdown export type: {destination.suffix}")

    return ConversionResult(
        source=source,
        destination=destination,
        markdown=markdown,
        warnings=warnings,
    )


def parse_markdown(markdown: str) -> list[MarkdownBlock]:
    lines = markdown.splitlines()
    blocks: list[MarkdownBlock] = []
    index = 0

    while index < len(lines):
        line = lines[index].rstrip()
        stripped = line.strip()

        if not stripped:
            index += 1
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading_match:
            blocks.append(HeadingBlock(level=len(heading_match.group(1)), text=heading_match.group(2).strip()))
            index += 1
            continue

        if _is_table_line(stripped) and index + 1 < len(lines) and _is_table_separator(lines[index + 1].strip()):
            table_lines = [stripped]
            index += 2
            while index < len(lines) and _is_table_line(lines[index].strip()):
                table_lines.append(lines[index].strip())
                index += 1
            blocks.append(TableBlock(rows=[_parse_table_row(item) for item in table_lines]))
            continue

        if _is_list_line(stripped):
            items: list[str] = []
            while index < len(lines):
                candidate = lines[index].strip()
                if not _is_list_line(candidate):
                    break
                items.append(re.sub(r"^([-*]|\d+\.)\s+", "", candidate).strip())
                index += 1
            blocks.append(ListBlock(items=items))
            continue

        paragraph_lines = [stripped]
        index += 1
        while index < len(lines):
            candidate = lines[index].strip()
            if not candidate:
                index += 1
                break
            if re.match(r"^(#{1,6})\s+", candidate) or _is_list_line(candidate):
                break
            if _is_table_line(candidate) and index + 1 < len(lines) and _is_table_separator(lines[index + 1].strip()):
                break
            paragraph_lines.append(candidate)
            index += 1
        blocks.append(ParagraphBlock(text=" ".join(paragraph_lines)))

    return blocks


def _write_docx(destination: Path, blocks: list[MarkdownBlock]) -> None:
    try:
        from docx import Document
    except ModuleNotFoundError as exc:
        raise RuntimeError("Markdown to Word conversion requires the `python-docx` package.") from exc

    document = Document()
    for block in blocks:
        if isinstance(block, HeadingBlock):
            document.add_heading(block.text, level=min(block.level, 4))
        elif isinstance(block, ParagraphBlock):
            document.add_paragraph(block.text)
        elif isinstance(block, ListBlock):
            for item in block.items:
                document.add_paragraph(item, style="List Bullet")
        elif isinstance(block, TableBlock) and block.rows:
            table = document.add_table(rows=len(block.rows), cols=max(len(row) for row in block.rows))
            table.style = "Table Grid"
            for row_index, row in enumerate(block.rows):
                for col_index, value in enumerate(row):
                    table.cell(row_index, col_index).text = value
    document.save(destination)


def _write_pptx(destination: Path, blocks: list[MarkdownBlock]) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError("Markdown to PowerPoint conversion requires the `python-pptx` package.") from exc

    presentation = Presentation()
    slides = _slides_from_blocks(blocks)
    title_block = blocks[0] if blocks and isinstance(blocks[0], HeadingBlock) and blocks[0].level == 1 else None

    if title_block:
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title_block.text
        title_slide.placeholders[1].text = "Converted from Markdown"

    for title, body_lines in slides or [("Markdown Export", ["Converted from Markdown."])]:
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for index, line in enumerate(body_lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0

    presentation.save(destination)


def _write_xlsx(destination: Path, blocks: list[MarkdownBlock]) -> None:
    try:
        from openpyxl import Workbook
    except ModuleNotFoundError as exc:
        raise RuntimeError("Markdown to Excel conversion requires the `openpyxl` package.") from exc

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Markdown"
    current_row = 1

    def next_sheet(title: str) -> tuple[object, int]:
        sheet_name = _safe_sheet_title(title)
        sheet = workbook.create_sheet(title=sheet_name)
        return sheet, 1

    for block in blocks:
        if isinstance(block, HeadingBlock) and block.level <= 2 and current_row > 1:
            worksheet, current_row = next_sheet(block.text)
            worksheet.cell(row=current_row, column=1, value=block.text)
            current_row += 2
            continue
        if isinstance(block, HeadingBlock):
            worksheet.cell(row=current_row, column=1, value=block.text)
            current_row += 2
        elif isinstance(block, ParagraphBlock):
            worksheet.cell(row=current_row, column=1, value=block.text)
            current_row += 2
        elif isinstance(block, ListBlock):
            for item in block.items:
                worksheet.cell(row=current_row, column=1, value=f"- {item}")
                current_row += 1
            current_row += 1
        elif isinstance(block, TableBlock):
            for row in block.rows:
                for col_index, value in enumerate(row, start=1):
                    worksheet.cell(row=current_row, column=col_index, value=value)
                current_row += 1
            current_row += 1

    workbook.save(destination)


def _write_pdf(destination: Path, blocks: list[MarkdownBlock]) -> None:
    lines = _pdf_lines_from_blocks(blocks)
    pdf = _build_basic_pdf(lines)
    destination.write_bytes(pdf)


def _slides_from_blocks(blocks: list[MarkdownBlock]) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title = "Markdown Export"
    current_lines: list[str] = []

    for block in blocks:
        if isinstance(block, HeadingBlock):
            if block.level == 1 and not slides and not current_lines:
                current_title = block.text
                continue
            if current_lines:
                slides.append((current_title, current_lines))
                current_lines = []
            current_title = block.text
            continue
        if isinstance(block, ParagraphBlock):
            current_lines.extend(textwrap.wrap(block.text, width=80) or [block.text])
        elif isinstance(block, ListBlock):
            current_lines.extend(f"- {item}" for item in block.items)
        elif isinstance(block, TableBlock):
            current_lines.extend(" | ".join(row) for row in block.rows)

    if current_lines or not slides:
        slides.append((current_title, current_lines or ["Converted from Markdown."]))
    return slides


def _pdf_lines_from_blocks(blocks: list[MarkdownBlock]) -> list[str]:
    lines: list[str] = []
    for block in blocks:
        if isinstance(block, HeadingBlock):
            lines.append(block.text.upper())
            lines.append("")
        elif isinstance(block, ParagraphBlock):
            lines.extend(textwrap.wrap(block.text, width=85) or [block.text])
            lines.append("")
        elif isinstance(block, ListBlock):
            for item in block.items:
                lines.extend(textwrap.wrap(f"- {item}", width=85) or [f"- {item}"])
            lines.append("")
        elif isinstance(block, TableBlock):
            for row in block.rows:
                lines.append(" | ".join(row))
            lines.append("")
    while lines and not lines[-1]:
        lines.pop()
    return lines or ["Converted from Markdown."]


def _build_basic_pdf(lines: list[str]) -> bytes:
    page_height = 792
    top_margin = 72
    bottom_margin = 72
    line_height = 16
    lines_per_page = max(1, (page_height - top_margin - bottom_margin) // line_height)
    pages = [lines[index:index + lines_per_page] for index in range(0, len(lines), lines_per_page)] or [[]]

    objects: list[bytes] = []
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")

    page_object_numbers: list[int] = []
    content_object_numbers: list[int] = []
    next_object = 4
    for _ in pages:
        page_object_numbers.append(next_object)
        content_object_numbers.append(next_object + 1)
        next_object += 2

    kids = " ".join(f"{number} 0 R" for number in page_object_numbers)
    objects.append(f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode("ascii"))
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    for page_number, page_lines in enumerate(pages):
        page_object = (
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Resources << /Font << /F1 3 0 R >> >> /Contents {content_object_numbers[page_number]} 0 R >>"
        ).encode("ascii")
        objects.append(page_object)

        content_lines = [
            "BT",
            "/F1 12 Tf",
            f"1 0 0 1 72 {page_height - top_margin} Tm",
            f"{line_height} TL",
        ]
        for line in page_lines:
            escaped = _escape_pdf_text(line)
            if escaped:
                content_lines.append(f"({escaped}) Tj")
            content_lines.append("T*")
        content_lines.append("ET")
        content_stream = "\n".join(content_lines).encode("latin-1", errors="replace")
        stream = (
            f"<< /Length {len(content_stream)} >>\nstream\n".encode("ascii")
            + content_stream
            + b"\nendstream"
        )
        objects.append(stream)

    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")

    xref_start = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_start}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(output)


def _escape_pdf_text(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _safe_sheet_title(value: str) -> str:
    cleaned = re.sub(r"[\\/*?:\[\]]", "-", value).strip()
    return (cleaned or "Markdown")[:31]


def _is_list_line(value: str) -> bool:
    return bool(re.match(r"^([-*]|\d+\.)\s+.+$", value))


def _is_table_line(value: str) -> bool:
    return value.startswith("|") and value.endswith("|") and value.count("|") >= 2


def _is_table_separator(value: str) -> bool:
    parts = [part.strip() for part in value.strip("|").split("|")]
    return bool(parts) and all(re.match(r"^:?-{3,}:?$", part) for part in parts)


def _parse_table_row(value: str) -> list[str]:
    return [cell.strip() for cell in value.strip("|").split("|")]
