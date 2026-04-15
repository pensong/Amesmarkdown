from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from mdconvert_app.assets import AssetManager
from mdconvert_app.markdown_utils import ensure_blank_line, markdown_table, normalize_whitespace
from mdconvert_app.models import ConversionResult


def convert_pptx(source: Path, destination: Path) -> ConversionResult:
    presentation = Presentation(source)
    assets = AssetManager(destination)
    lines: list[str] = [f"# {source.stem}", ""]

    for index, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"Slide {index}"
        lines.append(f"## {title}")
        lines.append("")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                relative = assets.save_bytes(image.blob, f".{image.ext}", f"slide-{index}-image")
                lines.append(f"![Slide {index} image]({relative})")
                lines.append("")
                continue
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([normalize_whitespace(cell.text) for cell in row.cells])
                if rows:
                    ensure_blank_line(lines)
                    lines.append(markdown_table(rows))
                    lines.append("")
                continue
            if getattr(shape, "has_text_frame", False):
                _append_text_frame(lines, shape.text_frame)

    markdown = "\n".join(_trim(lines)).strip() + "\n"
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(markdown, encoding="utf-8")
    return ConversionResult(source=source, destination=destination, markdown=markdown)


def _append_text_frame(lines: list[str], text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        text = normalize_whitespace("".join(run.text for run in paragraph.runs) or paragraph.text)
        if not text:
            continue
        level = getattr(paragraph, "level", 0) or 0
        bullet = "  " * level + "- "
        lines.append(f"{bullet}{text}")
    if lines and lines[-1] != "":
        lines.append("")


def _trim(lines: list[str]) -> list[str]:
    trimmed = list(lines)
    while trimmed and trimmed[-1] == "":
        trimmed.pop()
    return trimmed

