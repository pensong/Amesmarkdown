from pathlib import Path
from tempfile import TemporaryDirectory
import unittest
from unittest.mock import patch

from mdconvert_app.models import ConversionResult
from mdconvert_app.service import SUPPORTED_EXTENSIONS, convert_path

try:
    from docx import Document
except ModuleNotFoundError:
    Document = None

try:
    from openpyxl import load_workbook
except ModuleNotFoundError:
    load_workbook = None

try:
    from pptx import Presentation
except ModuleNotFoundError:
    Presentation = None


class ServiceTests(unittest.TestCase):
    def test_supported_extensions_include_pdf(self) -> None:
        self.assertIn(".pdf", SUPPORTED_EXTENSIONS)
        self.assertIn(".md", SUPPORTED_EXTENSIONS)

    def test_convert_path_routes_pdf_files(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "sample.pdf"
            source.write_bytes(b"%PDF-1.4\n")
            output_dir = Path(tmp) / "out"

            expected = ConversionResult(
                source=source.resolve(),
                destination=(output_dir / "sample.md").resolve(),
                markdown="# sample\n",
            )

            with patch("mdconvert_app.service._convert_file", return_value=expected) as convert_file:
                results = convert_path(source, output_dir)

        convert_file.assert_called_once_with(source.resolve(), (output_dir / "sample.md").resolve())
        self.assertEqual(results, [expected])

    @unittest.skipUnless(Document is not None, "python-docx is not installed")
    def test_convert_markdown_to_docx(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "sample.md"
            source.write_text("# Demo\n\nParagraph text.\n", encoding="utf-8")
            output_dir = Path(tmp) / "out"

            result = convert_path(source, output_dir, target_format="docx")[0]

            self.assertEqual(result.destination.suffix.lower(), ".docx")
            document = Document(result.destination)
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            self.assertIn("Demo", text)
            self.assertIn("Paragraph text.", text)

    @unittest.skipUnless(Presentation is not None, "python-pptx is not installed")
    def test_convert_markdown_to_pptx(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "slides.md"
            source.write_text("# Deck\n\n## Slide One\n\n- Alpha\n- Beta\n", encoding="utf-8")
            output_dir = Path(tmp) / "out"

            result = convert_path(source, output_dir, target_format="pptx")[0]

            self.assertEqual(result.destination.suffix.lower(), ".pptx")
            presentation = Presentation(result.destination)
            self.assertGreaterEqual(len(presentation.slides), 2)
            self.assertEqual(presentation.slides[0].shapes.title.text, "Deck")

    @unittest.skipUnless(load_workbook is not None, "openpyxl is not installed")
    def test_convert_markdown_to_xlsx(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "sheet.md"
            source.write_text("# Workbook\n\n| Name | Score |\n| --- | --- |\n| Alice | 10 |\n", encoding="utf-8")
            output_dir = Path(tmp) / "out"

            result = convert_path(source, output_dir, target_format="xlsx")[0]

            self.assertEqual(result.destination.suffix.lower(), ".xlsx")
            workbook = load_workbook(result.destination)
            sheet = workbook.active
            self.assertEqual(sheet["A1"].value, "Workbook")
            self.assertEqual(sheet["A3"].value, "Name")
            self.assertEqual(sheet["B4"].value, "10")

    def test_convert_markdown_to_pdf(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "report.md"
            source.write_text("# Report\n\nParagraph text.\n", encoding="utf-8")
            output_dir = Path(tmp) / "out"

            result = convert_path(source, output_dir, target_format="pdf")[0]

            self.assertEqual(result.destination.suffix.lower(), ".pdf")
            self.assertTrue(result.destination.read_bytes().startswith(b"%PDF-1.4"))

    def test_markdown_requires_target_format_for_directory_output(self) -> None:
        with TemporaryDirectory() as tmp:
            source = Path(tmp) / "sample.md"
            source.write_text("# Demo\n", encoding="utf-8")
            output_dir = Path(tmp) / "out"

            with self.assertRaisesRegex(ValueError, "Markdown input requires a target format"):
                convert_path(source, output_dir)


if __name__ == "__main__":
    unittest.main()
